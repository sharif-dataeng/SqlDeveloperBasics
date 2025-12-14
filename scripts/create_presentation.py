"""
Create a PPTX that embeds the SVG diagrams in docs\images.
Windows: install dependencies and run from repository root.

Requirements:
  python -m pip install python-pptx cairosvg Pillow

Run:
  cd c:\gitroot\SqlDeveloperBasicsX
  python .\scripts\create_presentation.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt
import os

ROOT = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
IMAGES_DIR = os.path.join(ROOT, "docs", "images")
OUT_DIR = os.path.join(ROOT, "ppt")
os.makedirs(OUT_DIR, exist_ok=True)
PPTX_PATH = os.path.join(OUT_DIR, "OLTP_OLAP_Normalization.pptx")

# Slide content (layman terms)
slides_content = [
    {
        "title": "OLTP — Everyday Transactions (Beginner)",
        "bullets": [
            "What: Systems that record everyday actions (buying, paying, updating).",
            "Example: When you buy something online, the order and payment are stored here.",
            "Key idea: Many small, fast operations that must be correct each time.",
            "Why it matters: Keeps live business data accurate (stock, balances).",
            "Typical design: uses many small tables to avoid duplicated data."
        ],
        "image": os.path.join(IMAGES_DIR, "oltp.png")
    },
    {
        "title": "OLAP — Analysis & Reports (Beginner)",
        "bullets": [
            "What: Systems used to answer big questions (sales by month, trends).",
            "Example: A dashboard that shows last year's sales by product.",
            "Key idea: Reads lots of historical data, runs slow but powerful queries.",
            "Why it matters: Helps businesses make decisions from past data.",
            "Typical design: stores summarized or flattened data for faster reports."
        ],
        "image": os.path.join(IMAGES_DIR, "olap.png")
    },
    {
        "title": "Normalized Data — Simple Explanation",
        "bullets": [
            "What: Data split into small tables so each fact is stored once.",
            "Analogy: One address book instead of copying addresses everywhere.",
            "Benefit: Easy and safe to update (no conflicting copies).",
            "Trade-off: Need to join tables to get full info (a bit slower to read).",
            "When used: Most transaction systems (OLTP)."
        ],
        "image": os.path.join(IMAGES_DIR, "normalization.png")
    },
    {
        "title": "Denormalized Data — Simple Explanation",
        "bullets": [
            "What: Data copied/flattened to make reading fast and simple.",
            "Analogy: A single printed report with everything on one page.",
            "Benefit: Very fast for reports and dashboards (fewer joins).",
            "Trade-off: Harder to keep updates consistent; uses more space.",
            "When used: Reporting systems (OLAP) and caches."
        ],
        "image": os.path.join(IMAGES_DIR, "denormalization.png")
    },
    {
        "title": "What is a Database? — Beginner",
        "bullets": [
            "What: A digital place that stores and organizes data safely.",
            "Analogy: A filing cabinet with rules for adding and finding papers.",
            "Types: Relational (tables), column stores (analytics), NoSQL (flexible).",
            "Why it matters: Lets apps remember users, orders, history securely.",
            "Quick tip: Use relational DBs for transactions, data warehouses for reports."
        ],
        "image": None
    }
]

prs = Presentation()
blank_layout = prs.slide_layouts[6]  # blank layout

for item in slides_content:
    slide = prs.slides.add_slide(blank_layout)
    # title box
    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.2), Inches(9), Inches(0.6))
    title_tf = title_box.text_frame
    title_tf.text = item["title"]
    title_tf.paragraphs[0].font.size = Pt(28)

    # body box for bullets
    body_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.0), Inches(8.5), Inches(4.5))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    first = True
    for b in item["bullets"]:
        if first:
            body_tf.text = b
            body_tf.paragraphs[0].font.size = Pt(18)
            first = False
        else:
            p = body_tf.add_paragraph()
            p.text = b
            p.level = 0
            p.font.size = Pt(18)

    # add image if exists (fits to slide)
    img_path = item.get("image")
    if img_path and os.path.exists(img_path):
        try:
            slide.shapes.add_picture(img_path, Inches(6.6), Inches(1.0), width=Inches(3.0))
        except Exception:
            # ignore image errors and continue
            pass

# save presentation
prs.save(PPTX_PATH)
print("Presentation created:", PPTX_PATH)